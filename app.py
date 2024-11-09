import os
import platform
from datetime import datetime, timedelta

import streamlit as st
import cv2
from PIL import Image
import imagehash
import pytesseract
import tempfile
import yt_dlp as youtube_dl
from pptx import Presentation
from pptx.util import Inches

if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = (
        r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    )


def convert_ms(timestamp_ms):
    # Convert milliseconds to seconds and create a timedelta object
    time_delta = timedelta(milliseconds=timestamp_ms)
    # Format as hours, minutes, seconds, and milliseconds
    readable_time = str(time_delta)
    if "." in readable_time:
        return "".join(readable_time.split(".")[:-1])
    return readable_time


def download_video(url, output_path="video.mp4", fast_mode=True):
    """Function to download the video from YouTube."""
    ydl_opts = {
        "format": "best"
        if fast_mode
        else "bestvideo",  # "best" is the fast but less quality.
        "outtmpl": output_path,
    }
    try:
        with youtube_dl.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        return output_path
    except Exception as e:
        st.error("Failed to download video. Please check the URL or try again later.")
        print(f"Error: {e}")
        return None


def extract_frames(video_path, frame_interval=30):
    """Function to extract frames from the video at specific intervals."""
    cap = cv2.VideoCapture(video_path)
    frame_id = 0
    frames = []

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break

        # Get the timestamp in milliseconds.
        timestamp = cap.get(cv2.CAP_PROP_POS_MSEC)

        # Save the frame and timestamp at the specified interval.
        if frame_id % frame_interval == 0:
            frames.append((frame, timestamp))

        frame_id += 1

    cap.release()
    return frames


def get_unique_slides(frames, hash_difference_threshold=5):
    """Function to get unique slides based on perceptual hashing."""
    hashes = []
    unique_frames = []

    for frame, timestamp in frames:
        pil_image = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
        hash_val = imagehash.phash(pil_image)

        # Only add frames that are sufficiently different from previous frames.
        if all(abs(hash_val - h) > hash_difference_threshold for h in hashes):
            unique_frames.append((frame, timestamp))
            hashes.append(hash_val)

    return unique_frames


def extract_text_from_slides(slides):
    """Function to extract text from each unique slide."""
    slide_texts = []

    for idx, (slide, timestamp) in enumerate(slides):
        pil_image = Image.fromarray(cv2.cvtColor(slide, cv2.COLOR_BGR2RGB))
        text = pytesseract.image_to_string(pil_image)

        # Convert to bytes for Streamlit compatibility.
        slide_image = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        pil_image.save(slide_image.name)

        slide_texts.append((slide_image.name, text.strip(), timestamp))

    return slide_texts


def create_ppt(slide_texts):
    """Function to create a PowerPoint presentation from slides."""
    prs = Presentation()

    for slide_filename, text, timestamp in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout.
        left = top = Inches(1)
        slide.shapes.add_picture(
            slide_filename, left, top, width=Inches(8.5), height=Inches(6)
        )
        txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8.5), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"At {convert_ms(timestamp)}\n{text}"

    # Save PPT to a temporary file.
    ppt_temp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(ppt_temp.name)
    return ppt_temp.name


def main():
    """Streamlit Interface"""
    st.title("YouTube Video Slide Extractor")
    st.write(
        "Enter a YouTube video link, and this app will extract unique slides along with any text on them."
    )

    youtube_url = st.text_input("YouTube video URL:", "")

    fast_mode = st.checkbox("Enable fast mode", value=True)
    if fast_mode:
        st.write("Fast video download but with less quality.")

    if youtube_url:
        if st.button("Process Video"):
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_video_path = os.path.join(temp_dir, "video.mp4")

                with st.spinner("Downloading video..."):
                    video_path = download_video(
                        youtube_url,
                        output_path=temp_video_path,
                        fast_mode=fast_mode,
                    )
                    if not video_path:
                        st.error("Video download failed. Please check the URL.")
                        return

                with st.spinner("Extracting frames..."):
                    frames = extract_frames(video_path)

                with st.spinner("Detecting unique slides..."):
                    unique_slides = get_unique_slides(frames)

                with st.spinner("Extracting text from slides..."):
                    slide_texts = extract_text_from_slides(unique_slides)

                ppt_filename = create_ppt(slide_texts)

                st.success("Processing complete!")
                st.write("Extracted Slides and Text:")

                for slide_filename, text, timestamp in slide_texts:
                    st.image(
                        slide_filename,
                        caption=f"Slide Image at {convert_ms(timestamp)}",
                        use_container_width=True,
                    )
                    st.write("Extracted Text:")
                    st.text(text)
                    st.write("---")

                # Download link for the PowerPoint file.
                with open(ppt_filename, "rb") as fd:
                    st.download_button(
                        "Download PowerPoint Presentation",
                        fd,
                        file_name="presentation.pptx",
                    )


if __name__ == "__main__":
    main()
